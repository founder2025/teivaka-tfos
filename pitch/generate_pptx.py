#!/usr/bin/env python3
"""Generate the Teivaka pitch deck as an editable .pptx (brand-themed, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- brand palette ----
CREAM=RGBColor(0xF8,0xF3,0xE9); PAPER=RGBColor(0xFF,0xFF,0xFF); SOIL=RGBColor(0x5C,0x40,0x33)
INK=RGBColor(0x2A,0x21,0x18); GREEN=RGBColor(0x6A,0xA8,0x4F); GREENDK=RGBColor(0x4F,0x8A,0x37)
AMBER=RGBColor(0xBF,0x90,0x00); MUTED=RGBColor(0x6A,0x5C,0x49); LINE=RGBColor(0xE2,0xD8,0xC3)
TINT=RGBColor(0xE8,0xF0,0xE0); DARK=RGBColor(0x2C,0x1A,0x0E); CREAMT=RGBColor(0xEC,0xE3,0xD1)
TEAL=RGBColor(0x4F,0x9D,0x87); WHITE=RGBColor(0xFF,0xFF,0xFF)
SERIF="Georgia"; SANS="Calibri"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=CREAM):
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=bg
    return s

def txt(s,l,t,w,h,text,size,color,bold=False,font=SANS,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,italic=False,sp=1.0,caps=False):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    lines=text.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        r=p.add_run(); r.text=ln.upper() if caps else ln
        f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return tb

def rect(s,l,t,w,h,fill=PAPER,line=LINE,line_w=1.0,rounded=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l),Inches(t),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False
    return shp

def badge(s,n,l=0.9,t=0.55):
    b=rect(s,l,t,0.62,0.62,fill=GREEN,line=None)
    tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(24); r.font.bold=True; r.font.name=SERIF; r.font.color.rgb=WHITE

def head(s,n,title):
    badge(s,n)
    txt(s,1.7,0.5,10.5,0.9,title,34,INK,bold=True,font=SERIF,anchor=MSO_ANCHOR.MIDDLE)

def eyebrow(s,text,l=0.9,t=0.6,color=GREENDK):
    txt(s,l,t,11,0.4,text,12,color,bold=True,font=SANS,caps=True)

def card(s,l,t,w,h,big,lbl,big_color=GREENDK,big_size=30):
    rect(s,l,t,w,h)
    txt(s,l+0.35,t+0.3,w-0.7,1.0,big,big_size,big_color,bold=True,font=SERIF)
    txt(s,l+0.35,t+1.25,w-0.7,h-1.4,lbl,13,SOIL,sp=1.05)

def foot(s,left,right,dark=False):
    c=CREAMT if dark else MUTED
    ln=rect(s,0.9,6.95,SW-1.8,0.012,fill=(RGBColor(0x55,0x40,0x33) if dark else LINE),line=None,rounded=False)
    txt(s,0.9,7.0,7.5,0.4,left,9.5,c)
    txt(s,SW-6.4,7.0,5.5,0.4,right,9.5,c,align=PP_ALIGN.RIGHT,italic=True)

# ---------- 1 COVER ----------
s=slide(DARK)
txt(s,0.9,1.3,11,1.0,"TEIVAKA",30,WHITE,bold=True,font=SERIF,caps=True)
txt(s,0.9,2.15,11,0.4,"Generate Wealth from Idle Lands · Fiji · Pre-seed",13,RGBColor(0x9E,0xD0,0x80),bold=True,caps=True)
txt(s,0.9,2.7,11.5,1.6,"Transform idle land into wealth.",52,WHITE,bold=True,font=SERIF,sp=0.95)
txt(s,0.9,4.5,9.5,1.0,"Fiji's first AI-powered agriculture operating system — making smallholder farmers visible, bankable, and profitable.",18,CREAMT,sp=1.2)
for i,(k,v) in enumerate([("Vision","Idle land → wealth"),("Mission","Every farmer prospers"),("Goal","Future of Pacific ag")]):
    x=0.9+i*3.3
    txt(s,x,5.7,3.1,0.3,k,10,RGBColor(0x9E,0xD0,0x80),bold=True,caps=True)
    txt(s,x,6.0,3.1,0.5,v,15,WHITE,bold=True,font=SERIF)
foot(s,"Uraia Koroi Kama (Cody), Founder · founder@teivaka.com","The Ask: FJD $75,000",dark=True)

# ---------- 2 PROBLEM ----------
s=slide(); head(s,1,"The problem")
txt(s,0.9,1.4,11.5,0.7,"Fiji has the land, the demand, and the farmers — yet families stay trapped in subsistence.",20,SOIL,italic=True,font=SERIF)
cw,ch=5.55,1.95; gx,gy=0.9,2.25; gap=0.45
cards=[("FJD 1.3B","imported in food every year — vs only ~FJD 150–250M in agri exports. A ~6:1 gap for food we can grow."),
       ("195,000 ha","already in farm holdings — under-producing, not unused. The land isn't missing; its potential is."),
       ("Farming by feel","plant, feed, harvest & sell on intuition → volatile yields and ~26–28% of harvest lost (est.) before sale."),
       ("No record","informal & undocumented — so banks, ministries & NGOs can't fund them. Productive farmers stay invisible.")]
for i,(b,l) in enumerate(cards):
    x=gx+(i%2)*(cw+gap); y=gy+(i//2)*(ch+0.3)
    card(s,x,y,cw,ch,b,l,big_size=26)
txt(s,0.9,6.45,11.5,0.4,"Sources: Fiji Bureau of Statistics (trade); Fiji National Agricultural Census; Fiji Development Bank. Post-harvest loss: sector estimate.",9.5,MUTED,italic=True)
foot(s,"Teivaka","The problem isn't land or market — it's the missing infrastructure to connect them.")

# ---------- 3 SOLUTION ----------
s=slide(); head(s,2,"The solution")
txt(s,0.9,1.4,11.5,0.7,"One login. One farmer record. Four pillars — and every pillar closes a money leak.",20,SOIL,italic=True,font=SERIF)
leaks=[("Low yields, thin margins — no field-tested system","TIS · AI",GREEN),
       ("Planting misses the import-price window","Classroom",AMBER),
       ("~26–28% of harvest lost before sale","TFOS",SOIL),
       ("Productive land idle — no record to unlock finance","TFOS",SOIL),
       ("Farmers sell blind — no live price signal","Community",TEAL),
       ("Hotels import what grows on the next farm over","Community",TEAL)]
ly=2.25
for lk,tag,col in leaks:
    txt(s,0.9,ly,6.0,0.4,lk,12.5,SOIL,anchor=MSO_ANCHOR.MIDDLE)
    tg=rect(s,6.9,ly+0.02,1.4,0.34,fill=col,line=None); tf=tg.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=tag; r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=WHITE
    ly+=0.62
box=rect(s,8.7,2.2,3.7,4.0,fill=SOIL,line=None)
txt(s,9.0,2.5,3.1,0.8,"The Teivaka Agriculture Ecosystem",18,WHITE,bold=True,font=SERIF,sp=1.0)
pill=[("Community","marketplace & live prices"),("TFOS","the record engine → bankable"),
      ("Classroom","practical local education"),("TIS","24/7 AI on WhatsApp (EN·Fijian·Hindi)")]
py=3.5
for nm,ds in pill:
    txt(s,9.0,py,3.2,0.3,nm,13,RGBColor(0x9E,0xD0,0x80),bold=True)
    txt(s,9.0,py+0.28,3.2,0.4,ds,11,CREAMT); py+=0.62
foot(s,"Teivaka","No feature without a dollar behind it.")

# ---------- 4 MARKET ----------
s=slide(); head(s,3,"Market opportunity")
txt(s,0.9,1.4,11.5,0.6,"The demand is already proven — by a FJD 1.3B import bill.",20,SOIL,italic=True,font=SERIF)
mc=[("83,000+","addressable smallholder farmers in Fiji — our beachhead."),
    ("14 nations","Pacific next — then global smallholder agriculture."),
    ("Farmers free","the organisations that benefit pay. 8 streams, no single dependency.")]
for i,(b,l) in enumerate(mc):
    card(s,0.9+i*4.0,2.2,3.7,2.1,b,l,big_size=28)
mix="Marketplace 30%   ·   Sponsored seats 20%   ·   Farmer subs 15%   ·   Agribusiness 10%   ·   Advertising 10%   ·   Education 5%   ·   Financial 5%   ·   Data 5%"
rect(s,0.9,4.7,11.5,1.0,fill=TINT,line=None)
txt(s,1.2,4.7,11.0,1.0,mix,14,GREENDK,bold=True,anchor=MSO_ANCHOR.MIDDLE,sp=1.3)
foot(s,"Teivaka","We don't charge the farmer — we charge the value that flows around the farmer.")

# ---------- 5 DELIVER & SCALE ----------
s=slide(); head(s,4,"How we deliver & scale")
txt(s,0.9,1.4,11.5,0.6,"Proven on our own farms first. Then scaled by software — at near-zero marginal cost.",20,SOIL,italic=True,font=SERIF)
ph=[("Phase 1 · Done","Build & prove","Full ecosystem built. 2 working farms returning $10–$12 per $1. Pre-revenue, founder-funded. Public site + TIS live on WhatsApp.",TINT,GREEN,GREENDK),
    ("Phase 2 · The $75k · 12 mo","Strengthen & launch","Harden the platform, 2→10 demo farms, onboard 250 farmers, complete legal, launch publicly. → ~FJD $40k ARR.",RGBColor(0xFB,0xF1,0xD8),AMBER,AMBER),
    ("Phase 3 · Forecast","Scale","Sponsored seats (NGO→Ministry) + marketplace become the engine. Pacific expansion. Sets up the seed round.",PAPER,LINE,MUTED)]
for i,(tag,title,body,fill,ln,tagc) in enumerate(ph):
    x=0.9+i*4.0; rect(s,x,2.3,3.7,3.3,fill=fill,line=ln,line_w=1.5)
    txt(s,x+0.3,2.55,3.2,0.4,tag,11,tagc,bold=True,caps=True)
    txt(s,x+0.3,3.0,3.2,0.6,title,18,INK,bold=True,font=SERIF)
    txt(s,x+0.3,3.7,3.2,1.7,body,12,SOIL,sp=1.1)
foot(s,"Teivaka","QA = cited knowledge base (no invented agronomy) + audit-anchored records.")

# ---------- 6 ECONOMICS ----------
s=slide(); head(s,5,"The economics")
txt(s,0.9,1.4,11.5,0.6,"We built all of this on almost nothing — funded by the farms, not investors.",20,SOIL,italic=True,font=SERIF)
ec=[("$850/mo","total software burn — funded by the pilot farms. Zero outside capital."),
    ("$10–$12","returned for every $1 spent on the farms. We know our numbers."),
    ("~FJD 40k","target ARR in 12 months — subs + agribusiness + ads + marketplace + sponsored seats.")]
for i,(b,l) in enumerate(ec):
    card(s,0.9+i*4.0,2.2,3.7,2.0,b,l,big_size=30)
rect(s,0.9,4.6,11.5,1.5,fill=SOIL,line=None)
txt(s,1.3,4.6,10.7,1.5,"Breakeven is one institutional deal away. One NGO sponsoring 500 farmers (FJD $10k/yr) covers our entire annual burn. One Ministry package (FJD $75k/yr) covers it seven times over.",18,CREAM,font=SERIF,anchor=MSO_ANCHOR.MIDDLE,sp=1.15)
foot(s,"Teivaka","Pre-revenue — by design, not by accident.")

# ---------- 7 MOAT ----------
s=slide(); head(s,6,"Why a neighbour can't copy us tomorrow")
mo=[("The data moat compounds","years of verified farm records can't be faked or rushed. The earlier we start, the bigger the lead."),
    ("Cited, local AI","a Fiji-specific knowledge base, in-language — generic AI can't match it, and it never invents agronomy."),
    ("Ecosystem lock-in","one login across market + records + learning + AI = real switching cost."),
    ("Distribution + credibility","WhatsApp-native, low-literacy accessible — and built by farmers, on profitable farms.")]
for i,(b,l) in enumerate(mo):
    x=0.9+(i%2)*5.95; y=1.7+(i//2)*2.4
    rect(s,x,y,5.5,2.1)
    txt(s,x+0.35,y+0.3,4.9,0.7,b,18,GREENDK,bold=True,font=SERIF)
    txt(s,x+0.35,y+1.1,4.9,0.9,l,13,SOIL,sp=1.1)
foot(s,"Teivaka","We give the software away — and that's exactly why it's defensible.")

# ---------- 8 TEAM ----------
s=slide(); head(s,7,"The team")
txt(s,0.9,1.4,11.5,0.6,"Four Fijians. We run the farms. We build the platform. Same people.",20,SOIL,italic=True,font=SERIF)
team=[("UK","Uraia Koroi Kama","Founder & Director","Strategy, partnerships, platform architecture. Left a Science Degree; seven years farming. Knows every pain point — and builds the platform."),
      ("TK","Taniela Kama","Field Ops & Nursery","Nursery, land prep, farm layout. 20+ years of Fiji farming knowledge. Ready before planting."),
      ("LW","Laisenia Waqa","Crop Production & Harvest","Full crop cycle — pest, weed, nutrient. Quality control. Market-grade yield."),
      ("KW","Kinisimere Wati","Finance & Operations","Bookkeeping, payroll, revenue, invoices. Every figure tight and transparent.")]
cw=2.85; gx=0.55
for i,(ini,nm,role,bio) in enumerate(team):
    x=gx+i*3.05; rect(s,x,2.2,cw,3.5)
    av=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+cw/2-0.45),Inches(2.45),Inches(0.9),Inches(0.9))
    av.fill.solid(); av.fill.fore_color.rgb=GREEN; av.line.fill.background(); av.shadow.inherit=False
    tf=av.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=ini; r.font.size=Pt(20); r.font.bold=True; r.font.name=SERIF; r.font.color.rgb=WHITE
    txt(s,x+0.15,3.5,cw-0.3,0.5,nm,15,INK,bold=True,font=SERIF,align=PP_ALIGN.CENTER)
    txt(s,x+0.1,4.05,cw-0.2,0.4,role,9.5,GREENDK,bold=True,align=PP_ALIGN.CENTER,caps=True)
    txt(s,x+0.2,4.5,cw-0.4,1.1,bio,10.5,SOIL,align=PP_ALIGN.CENTER,sp=1.05)
txt(s,0.9,5.95,11.5,0.5,"100% Fijian. 100% local ownership. Female finance lead by design.",15,MUTED,italic=True,font=SERIF,align=PP_ALIGN.CENTER)
foot(s,"Teivaka","Founder-operators — we run what we build.")

# ---------- 9 THE ASK ----------
s=slide(); head(s,8,"The ask")
txt(s,0.9,1.5,4.5,0.4,"PRE-SEED",13,GREENDK,bold=True,caps=True)
txt(s,0.85,1.9,5.0,1.3,"FJD $75,000",48,GREENDK,bold=True,font=SERIF)
txt(s,0.9,3.3,5.2,1.0,"To go from a founder-funded pilot to a commercially ready ecosystem — and launch publicly in 12 months.",15,SOIL,sp=1.2)
txt(s,0.9,4.6,5.4,1.3,"12-mo: 250 registered · 150 active · 50 paying · 20 agribusiness · 10 farms · ~FJD $40k ARR",13,MUTED,bold=True,sp=1.25)
rows=[("Product strengthening & commercial readiness","22,500"),
      ("Pilot network & data validation (2→10 farms)","15,000"),
      ("Farmer acquisition & onboarding (250)","18,750"),
      ("Operations & infrastructure (12 mo)","7,500"),
      ("Legal, compliance & governance","5,250"),
      ("Launch reserve (buffer)","6,000"),
      ("Total","FJD 75,000")]
tbl=s.shapes.add_table(len(rows),2,Inches(6.6),Inches(1.7),Inches(5.8),Inches(4.4)).table
tbl.columns[0].width=Inches(4.4); tbl.columns[1].width=Inches(1.4)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        c=tbl.cell(ri,ci); c.fill.solid(); c.fill.fore_color.rgb=(TINT if ri==len(rows)-1 else PAPER)
        c.margin_left=Inches(0.12); c.margin_top=Inches(0.04); c.margin_bottom=Inches(0.04)
        tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
        p.alignment=PP_ALIGN.RIGHT if ci==1 else PP_ALIGN.LEFT
        r=p.add_run(); r.text=val; r.font.size=Pt(12); r.font.name=SANS
        r.font.bold=(ri==len(rows)-1 or ci==1); r.font.color.rgb=INK if ci==1 else SOIL
foot(s,"Teivaka","Aggressive enough to matter. Conservative enough to defend.")

# ---------- 10 IMPACT ----------
s=slide(); head(s,9,"The impact")
txt(s,0.9,1.4,11.5,0.6,"When a farmer becomes visible, everything downstream changes.",20,SOIL,italic=True,font=SERIF)
imp=[("Economic",["Idle land activated into income","Records make farmers bankable — credit flows","Local supply chips at the $1.3B import bill"]),
     ("Social",["Knowledge & markets reach every farmer — in their language","100% Fijian, 100% locally owned","Female finance lead by design"]),
     ("Environmental",["Less post-harvest waste (~26–28% today)","Better-timed inputs, healthier soil","Productive use of existing land — no new clearing"])]
for i,(t,items) in enumerate(imp):
    x=0.9+i*4.0; rect(s,x,2.2,3.7,2.9)
    tg=rect(s,x+0.3,2.45,1.7,0.4,fill=TINT,line=None); tf=tg.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=GREENDK
    bt=s.shapes.add_textbox(Inches(x+0.3),Inches(3.0),Inches(3.1),Inches(2.0)); bf=bt.text_frame; bf.word_wrap=True
    for j,it in enumerate(items):
        p=bf.paragraphs[0] if j==0 else bf.add_paragraph(); p.line_spacing=1.05
        r=p.add_run(); r.text="• "+it; r.font.size=Pt(11.5); r.font.color.rgb=SOIL; r.font.name=SANS
rect(s,0.9,5.35,11.5,1.1,fill=SOIL,line=None)
txt(s,1.3,5.35,10.7,1.1,"The land, the demand, and the farmers are already here. Teivaka is the operating system that finally connects them.",18,CREAM,font=SERIF,anchor=MSO_ANCHOR.MIDDLE,sp=1.1)
foot(s,"Uraia Koroi Kama (Cody) · founder@teivaka.com · teivaka.com","Build the future of Pacific agriculture.")

prs.save("pitch/teivaka_pitch_deck.pptx")
print("saved pitch/teivaka_pitch_deck.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
